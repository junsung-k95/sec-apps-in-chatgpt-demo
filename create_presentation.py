"""Samsung Galaxy MCP Server - Presentation Generator (Samsung Blue Style)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Samsung Brand Colors ──
SAMSUNG_BLUE = RGBColor(0x14, 0x28, 0xA0)      # Primary brand blue
DARK_BLUE = RGBColor(0x0A, 0x1A, 0x5C)          # Dark navy
MID_BLUE = RGBColor(0x1E, 0x3A, 0xB8)           # Mid blue
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)         # Accent light blue
SKY_BLUE = RGBColor(0xD6, 0xE8, 0xF7)           # Very light blue bg
ICE_BLUE = RGBColor(0xEE, 0xF4, 0xFB)           # Ice background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x00, 0x9E, 0x60)
ORANGE = RGBColor(0xF0, 0x7E, 0x13)
CARE_PLUS_BLUE = RGBColor(0x00, 0x6E, 0xBE)
CLUB_PURPLE = RGBColor(0x5C, 0x2D, 0x91)
TRADEIN_GREEN = RGBColor(0x00, 0x7A, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helper Functions ──

def add_bg_rect(slide, color, left=0, top=0, width=None, height=None):
    """Add a solid color rectangle as background."""
    w = width or W
    h = height or H
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="맑은 고딕", line_spacing=1.2, anchor=MSO_ANCHOR.TOP):
    """Add a text box with a single run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_multi_para(slide, left, top, width, height, paragraphs, font_size=13,
                   color=BLACK, font_name="맑은 고딕", line_spacing=1.3,
                   anchor=MSO_ANCHOR.TOP):
    """
    Add text box with multiple paragraphs.
    paragraphs: list of (text, bold, color, font_size_override, alignment)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, para_data in enumerate(paragraphs):
        if isinstance(para_data, str):
            para_data = (para_data, False, color, font_size, PP_ALIGN.LEFT)
        text, bld, clr, fs, align = para_data

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        p.line_spacing = Pt(fs * line_spacing)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.color.rgb = clr
        run.font.bold = bld
        run.font.name = font_name
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, radius=Inches(0.15)):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, items, title_color=SAMSUNG_BLUE,
             bg_color=WHITE, border_color=SKY_BLUE, title_size=14, item_size=11):
    """Add a card-style box with title and bullet items."""
    add_rounded_rect(slide, left, top, width, height, bg_color, border_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.4),
                 title, font_size=title_size, color=title_color, bold=True)
    paras = []
    for item in items:
        paras.append((f"• {item}", False, DARK_GRAY, item_size, PP_ALIGN.LEFT))
    if paras:
        add_multi_para(slide, left + Inches(0.2), top + Inches(0.5),
                       width - Inches(0.4), height - Inches(0.6), paras,
                       font_size=item_size, color=DARK_GRAY)


def add_slide_number(slide, num, total):
    """Add slide number at bottom right."""
    add_text_box(slide, W - Inches(1.2), H - Inches(0.45), Inches(1), Inches(0.35),
                 f"{num} / {total}", font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_header_bar(slide, title_text):
    """Add a dark blue header bar across the top."""
    add_bg_rect(slide, DARK_BLUE, 0, 0, W, Inches(1.05))
    add_text_box(slide, Inches(0.7), Inches(0.2), Inches(10), Inches(0.65),
                 title_text, font_size=26, color=WHITE, bold=True)
    # Thin accent line
    add_bg_rect(slide, LIGHT_BLUE, 0, Inches(1.05), W, Inches(0.04))


def add_samsung_footer(slide):
    """Add Samsung-style footer line."""
    add_bg_rect(slide, SKY_BLUE, 0, H - Inches(0.5), W, Inches(0.5))
    add_text_box(slide, Inches(0.7), H - Inches(0.43), Inches(4), Inches(0.35),
                 "Samsung Galaxy Services — MCP Server", font_size=9, color=MID_BLUE, bold=False)


TOTAL_SLIDES = 14


# ══════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg_rect(slide, DARK_BLUE)

# Decorative gradient-like strips
add_bg_rect(slide, MID_BLUE, Inches(0), Inches(3.0), W, Inches(0.06))
add_bg_rect(slide, LIGHT_BLUE, Inches(0), Inches(5.2), W, Inches(0.04))

# Samsung logo placeholder text
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.5),
             "SAMSUNG", font_size=20, color=WHITE, bold=True, font_name="Arial")

# Main title
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
             "Samsung Galaxy Services", font_size=44, color=WHITE, bold=True, font_name="맑은 고딕")
add_text_box(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(0.9),
             "ChatGPT 연동 MCP 서버 — 기능 소개 및 기술 요건", font_size=22, color=SKY_BLUE, bold=False)

# Three service badges
badge_data = [
    ("Samsung Care+", "기기 보험 서비스", CARE_PLUS_BLUE),
    ("New 갤럭시 AI 구독클럽", "구독형 업그레이드", CLUB_PURPLE),
    ("보상판매 (Trade-in)", "기기 보상가 견적", TRADEIN_GREEN),
]
for i, (name, desc, badge_color) in enumerate(badge_data):
    bx = Inches(0.8) + Inches(i * 3.8)
    by = Inches(4.6)
    add_rounded_rect(slide, bx, by, Inches(3.4), Inches(1.1), badge_color)
    add_text_box(slide, bx + Inches(0.25), by + Inches(0.15), Inches(2.9), Inches(0.45),
                 name, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, bx + Inches(0.25), by + Inches(0.6), Inches(2.9), Inches(0.4),
                 desc, font_size=12, color=RGBColor(0xDD, 0xDD, 0xFF))

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(5), Inches(0.4),
             "2026.03", font_size=14, color=RGBColor(0x88, 0xAA, 0xDD))
add_slide_number(slide, 1, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2: Project Overview
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "프로젝트 개요")
add_samsung_footer(slide)

# Subtitle
add_text_box(slide, Inches(0.7), Inches(1.3), Inches(11), Inches(0.5),
             "ChatGPT 안에서 삼성 Galaxy 서비스를 바로 이용할 수 있다면?",
             font_size=20, color=MID_BLUE, bold=True)

# Concept explanation
add_text_box(slide, Inches(0.7), Inches(2.0), Inches(11.5), Inches(0.8),
             "이 프로젝트는 ChatGPT 대화 안에서 삼성전자의 Galaxy 서비스 3가지를\n"
             "대화형으로 안내하고, 실시간 위젯(화면)으로 정보를 보여주는 시스템입니다.",
             font_size=14, color=DARK_GRAY)

# Flow diagram
# User → ChatGPT → MCP Server → Widget
flow_boxes = [
    ("👤 사용자", "자연어로 질문", RGBColor(0x4A, 0x4A, 0x4A)),
    ("🤖 ChatGPT", "질문 이해 + 도구 선택", MID_BLUE),
    ("⚙️ MCP 서버", "정보 조회 + 비즈니스 로직", SAMSUNG_BLUE),
    ("📊 위젯 + 텍스트", "사용자에게 결과 표시", CARE_PLUS_BLUE),
]
for i, (title, desc, color) in enumerate(flow_boxes):
    fx = Inches(0.7) + Inches(i * 3.1)
    fy = Inches(3.1)
    add_rounded_rect(slide, fx, fy, Inches(2.7), Inches(1.3), color)
    add_text_box(slide, fx + Inches(0.15), fy + Inches(0.15), Inches(2.4), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, fx + Inches(0.15), fy + Inches(0.7), Inches(2.4), Inches(0.45),
                 desc, font_size=11, color=RGBColor(0xDD, 0xEE, 0xFF), alignment=PP_ALIGN.CENTER)
    if i < 3:
        # Arrow
        add_text_box(slide, fx + Inches(2.55), fy + Inches(0.35), Inches(0.6), Inches(0.5),
                     "→", font_size=24, color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Three service cards
svc_cards = [
    ("Samsung Care+", ["기기 보험 가입 상담", "Late Enrollment (사진 검사)", "Basic / Premium 비교"], CARE_PLUS_BLUE),
    ("New 갤럭시 AI 구독클럽", ["12/24/36개월 플랜 안내", "구독 vs 일반구매 비용 비교", "가입 절차 안내"], CLUB_PURPLE),
    ("보상판매 (Trade-in)", ["기기 보상가 간편 조회", "정식 견적 (상태별 가격 산정)", "사진 기반 정밀 재산정"], TRADEIN_GREEN),
]
for i, (name, items, accent) in enumerate(svc_cards):
    cx = Inches(0.7) + Inches(i * 4.1)
    cy = Inches(4.8)
    add_rounded_rect(slide, cx, cy, Inches(3.7), Inches(2.0), ICE_BLUE, accent)
    add_text_box(slide, cx + Inches(0.2), cy + Inches(0.12), Inches(3.3), Inches(0.4),
                 name, font_size=14, color=accent, bold=True)
    for j, item in enumerate(items):
        add_text_box(slide, cx + Inches(0.25), cy + Inches(0.55 + j * 0.38), Inches(3.2), Inches(0.35),
                     f"• {item}", font_size=11, color=DARK_GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3: Scenario 1 — Care+ Late Enrollment
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "시나리오 1 — Samsung Care+ Late Enrollment")
add_samsung_footer(slide)

# Scenario context
add_rounded_rect(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.55), ICE_BLUE, LIGHT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.45),
             "💬 상황: Galaxy S24 Ultra를 작년 10월에 산 고객이 지금도 Care+에 가입할 수 있는지 문의",
             font_size=13, color=MID_BLUE, bold=True)

# Turn 1
y = Inches(2.0)
add_text_box(slide, Inches(0.7), y, Inches(3), Inches(0.35),
             "턴 1  가입 가능 여부 확인", font_size=13, color=SAMSUNG_BLUE, bold=True)
y += Inches(0.35)

# User bubble
add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.85), y + Inches(0.08), Inches(5.2), Inches(0.35),
             '👤 "S24 Ultra를 작년 10월에 샀는데, Care+에 지금도 가입할 수 있나요?"',
             font_size=11, color=BLACK)

# MCP call box
add_rounded_rect(slide, Inches(6.5), y - Inches(0.05), Inches(6.1), Inches(2.1), RGBColor(0xF0, 0xF5, 0xFF), CARE_PLUS_BLUE)
add_text_box(slide, Inches(6.7), y, Inches(3), Inches(0.3),
             '⚙️ Care+ 플랜 조회 및 가입 가능 여부 확인', font_size=12, color=CARE_PLUS_BLUE, bold=True)
mcp_paras = [
    ("입력: 기기=S24 Ultra, 구매일=2025-10-01", False, MID_GRAY, 10, PP_ALIGN.LEFT),
    ("", False, MID_GRAY, 6, PP_ALIGN.LEFT),
    ("출력:", True, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • Care+ 플랜 목록 (Basic $11.99/월, Premium $17.99/월)", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • 가입 상태: Late Enrollment 가능 (60일~1년)", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • Late Enrollment 수수료: $49", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • Vision 검사 필요: Yes", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
]
add_multi_para(slide, Inches(6.7), y + Inches(0.32), Inches(5.7), Inches(1.7), mcp_paras)

# AI response
y += Inches(0.6)
add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.65), SKY_BLUE)
add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(5.2), Inches(0.55),
             '🤖 "Late Enrollment로 가입 가능합니다. 기기 사진을 올려주세요"\n'
             '    [Care+ 위젯: Basic/Premium 비교 + Late Enrollment 배너]',
             font_size=10, color=DARK_BLUE)

# Turn 2
y += Inches(0.95)
add_text_box(slide, Inches(0.7), y, Inches(3), Inches(0.35),
             "턴 2  사진 기반 기기 상태 확인 (AI Vision)", font_size=13, color=SAMSUNG_BLUE, bold=True)
y += Inches(0.35)

add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.85), y + Inches(0.08), Inches(5.2), Inches(0.35),
             '👤 [Galaxy S24 Ultra 앞면/뒷면 사진 2장 업로드]',
             font_size=11, color=BLACK)

# MCP call box 2
add_rounded_rect(slide, Inches(6.5), y - Inches(0.05), Inches(6.1), Inches(2.0), RGBColor(0xF0, 0xF5, 0xFF), CARE_PLUS_BLUE)
add_text_box(slide, Inches(6.7), y, Inches(4), Inches(0.3),
             '⚙️ 사진 기반 기기 상태 검사 (Late Enrollment)', font_size=12, color=CARE_PLUS_BLUE, bold=True)
mcp_paras2 = [
    ("입력: 기기=S24 Ultra, 구매일=2025-10-01", False, MID_GRAY, 10, PP_ALIGN.LEFT),
    ("       Vision 결과: 화면=미세스크래치, 외관=미세흔적, 카메라=깨끗", False, MID_GRAY, 10, PP_ALIGN.LEFT),
    ("", False, MID_GRAY, 6, PP_ALIGN.LEFT),
    ("출력:", True, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • 적격 여부: ✅ Late Enrollment 승인", False, GREEN, 10, PP_ALIGN.LEFT),
    ("  • 항목별: 화면 ✅  외관 ✅  카메라 ✅", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • 추천 플랜: Care+ Premium", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
]
add_multi_para(slide, Inches(6.7), y + Inches(0.32), Inches(5.7), Inches(1.5), mcp_paras2)

y += Inches(0.6)
add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.65), SKY_BLUE)
add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(5.2), Inches(0.55),
             '🤖 "기기 상태 확인 완료! Care+ 가입이 가능합니다."\n'
             '    [위젯: 상태 🟢🟢🟢 + "Late Enrollment 승인" 녹색 배너]',
             font_size=10, color=DARK_BLUE)

# Turn 3-4 (compact)
y += Inches(0.8)
add_rounded_rect(slide, Inches(0.7), y, Inches(11.9), Inches(0.55), LIGHT_GRAY)
add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(11.5), Inches(0.45),
             '턴 3-4  👤 "Premium이랑 Basic 차이가 뭐예요? 분실도 보장되나요?" → 🤖 상세 비교 표시 → 👤 "Premium으로 가입할게요" → 🤖 CTA "Care+ Premium 가입하기"',
             font_size=11, color=DARK_GRAY)

add_slide_number(slide, 3, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4: Scenario 2 — Galaxy Club (Updated)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "시나리오 2 — New 갤럭시 AI 구독클럽 상담")
add_samsung_footer(slide)

add_rounded_rect(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.55), ICE_BLUE, LIGHT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.45),
             "💬 상황: 갤럭시 AI 구독클럽이 뭔지 궁금하고, 가입 시 어떤 혜택이 있는지 알고 싶은 고객",
             font_size=13, color=MID_BLUE, bold=True)

# Turn 1
y = Inches(2.0)
add_text_box(slide, Inches(0.7), y, Inches(4), Inches(0.35),
             "턴 1  구독클럽 소개", font_size=13, color=SAMSUNG_BLUE, bold=True)
y += Inches(0.35)

add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.85), y + Inches(0.08), Inches(5.2), Inches(0.35),
             '👤 "갤럭시 AI 구독클럽이 뭔가요?"',
             font_size=11, color=BLACK)

add_rounded_rect(slide, Inches(6.5), y - Inches(0.05), Inches(6.1), Inches(2.3), RGBColor(0xF5, 0xF0, 0xFF), CLUB_PURPLE)
add_text_box(slide, Inches(6.7), y, Inches(4), Inches(0.3),
             '⚙️ Galaxy Club 플랜별 가격·혜택 조회', font_size=12, color=CLUB_PURPLE, bold=True)
mcp_p = [
    ("입력: 플랜=all, FAQ=true", False, MID_GRAY, 10, PP_ALIGN.LEFT),
    ("", False, MID_GRAY, 6, PP_ALIGN.LEFT),
    ("출력:", True, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • 플랜 3종:", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("    - 12개월형: ₩6,900/월, 반납 시 기준가 50% 현금 지급", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("    - 24개월형: ₩6,900/월, 반납 시 40% 현금 지급", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("    - 36개월형: ₩8,900/월, 반납 시 25% + 분실보장", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • Care+ 포함, 액세서리 쿠폰, 가입절차 3단계", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • 라이프사이클 타임라인", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
]
add_multi_para(slide, Inches(6.7), y + Inches(0.32), Inches(5.7), Inches(1.8), mcp_p)

y += Inches(0.6)
add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.75), SKY_BLUE)
add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(5.2), Inches(0.65),
             '🤖 "삼성닷컴에서 기기 구매 시 가입하는 부가 구독 프로그램입니다.\n'
             '     Care+, 잔존가 보장, 액세서리 쿠폰이 포함됩니다.\n'
             '     관심 기기가 있으면 비용 구조를 분석해 드릴까요?"  [구독클럽 위젯]',
             font_size=10, color=DARK_BLUE)

# Turn 2
y += Inches(1.2)
add_text_box(slide, Inches(0.7), y, Inches(5), Inches(0.35),
             "턴 2  비용 구조 분석 (가입 시 vs 미가입 시)", font_size=13, color=SAMSUNG_BLUE, bold=True)
y += Inches(0.35)

add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.85), y + Inches(0.08), Inches(5.2), Inches(0.35),
             '👤 "Galaxy S26 Ultra로 비교해주세요"',
             font_size=11, color=BLACK)

add_rounded_rect(slide, Inches(6.5), y - Inches(0.05), Inches(6.1), Inches(2.2), RGBColor(0xF5, 0xF0, 0xFF), CLUB_PURPLE)
add_text_box(slide, Inches(6.7), y, Inches(4), Inches(0.3),
             '⚙️ 구독 가입 시 vs 미가입 시 비용 비교', font_size=12, color=CLUB_PURPLE, bold=True)
mcp_p2 = [
    ("입력: 기기=Galaxy S26 Ultra, 플랜=12개월형", False, MID_GRAY, 10, PP_ALIGN.LEFT),
    ("", False, MID_GRAY, 6, PP_ALIGN.LEFT),
    ("출력:", True, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  ▸ 가입 시: 기기 1,753,650원(공통) + 이용료 82,800원", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("    Care+ 포함, 반납 시 잔존가 898,700원 현금 지급", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("    → 반납 후 실질 부담: 937,750원", False, CLUB_PURPLE, 10, PP_ALIGN.LEFT),
    ("  ▸ 미가입: 기기값 + Care+ 별도 290,250원/년", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("    중고 판매 가격 보장 없음", False, MID_GRAY, 10, PP_ALIGN.LEFT),
]
add_multi_para(slide, Inches(6.7), y + Inches(0.32), Inches(5.7), Inches(1.7), mcp_p2)

y += Inches(0.6)
add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.55), SKY_BLUE)
add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(5.2), Inches(0.45),
             '🤖 [비교 위젯: 가입 시 실질부담 937,750원 vs 미가입 시 보장 없음]',
             font_size=10, color=DARK_BLUE)

# Turn 3 compact
y += Inches(0.75)
add_rounded_rect(slide, Inches(0.7), y, Inches(11.9), Inches(0.55), LIGHT_GRAY)
add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(11.5), Inches(0.45),
             '턴 3  👤 "가입하고 싶어요" → 🤖 3단계 안내: 삼성닷컴 구매+신청 → IMEI 등록 → 반납+잔존가 수령 + CTA "가입하기"',
             font_size=11, color=DARK_GRAY)

add_slide_number(slide, 4, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5: Scenario 3 — Trade-in (Updated)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "시나리오 3 — Trade-in 보상판매 견적")
add_samsung_footer(slide)

add_rounded_rect(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.55), ICE_BLUE, LIGHT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.45),
             "💬 상황: 기존 기기의 보상판매 가격을 알고 싶고, 사진으로 더 정확한 견적을 받고 싶은 고객",
             font_size=13, color=MID_BLUE, bold=True)

# Turn 1 — 대화로 정보 수집
y = Inches(2.0)
add_text_box(slide, Inches(0.7), y, Inches(4), Inches(0.35),
             "턴 1  보상판매 문의", font_size=13, color=SAMSUNG_BLUE, bold=True)
y += Inches(0.35)

add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.85), y + Inches(0.08), Inches(5.2), Inches(0.35),
             '👤 "보상판매 얼마 받을 수 있는지 알아보고 싶어요"',
             font_size=11, color=BLACK)

# No MCP call yet — AI asks follow-up
y += Inches(0.6)
add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.55), SKY_BLUE)
add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(5.2), Inches(0.45),
             '🤖 "어떤 기기를 보상판매하시려고요?\n     그리고 어느 통신사를 사용하고 계신가요?"',
             font_size=10, color=DARK_BLUE)

# Right side — note
add_rounded_rect(slide, Inches(6.5), y - Inches(0.65), Inches(6.1), Inches(1.1), RGBColor(0xF0, 0xFF, 0xF0), TRADEIN_GREEN)
add_text_box(slide, Inches(6.7), y - Inches(0.6), Inches(5.7), Inches(0.9),
             '이 시점에서는 MCP 서버 호출 없이\nChatGPT가 대화로 필요한 정보(기기명, 통신사)를\n수집합니다. 가이드라인에 따라 자연스럽게 질문.',
             font_size=10, color=DARK_GRAY)

# Turn 2 — 기기+통신사로 기본 견적
y += Inches(0.85)
add_text_box(slide, Inches(0.7), y, Inches(4), Inches(0.35),
             "턴 2  기기+통신사 기반 기본 견적", font_size=13, color=SAMSUNG_BLUE, bold=True)
y += Inches(0.35)

add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.85), y + Inches(0.08), Inches(5.2), Inches(0.35),
             '👤 "Galaxy S23 Ultra, SKT입니다"',
             font_size=11, color=BLACK)

add_rounded_rect(slide, Inches(6.5), y - Inches(0.05), Inches(6.1), Inches(1.8), RGBColor(0xF0, 0xFF, 0xF0), TRADEIN_GREEN)
add_text_box(slide, Inches(6.7), y, Inches(4), Inches(0.3),
             '⚙️ 보상판매 정식 견적 시작', font_size=12, color=TRADEIN_GREEN, bold=True)
mcp_t2 = [
    ("입력: 기기=Galaxy S23 Ultra, 통신사=SKT", False, MID_GRAY, 10, PP_ALIGN.LEFT),
    ("", False, MID_GRAY, 6, PP_ALIGN.LEFT),
    ("출력:", True, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • 견적 ID 발급", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • 기본 견적가 산출 (모델+통신사 기반)", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • 가격 산출 내역 (기준가, 통신사 조정 등)", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
]
add_multi_para(slide, Inches(6.7), y + Inches(0.32), Inches(5.7), Inches(1.3), mcp_t2)

y += Inches(0.6)
add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.55), SKY_BLUE)
add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(5.2), Inches(0.45),
             '🤖 [Trade-in 위젯: 초기 견적 표시]\n     "더 정확한 견적을 위해 기기 사진을 올려주세요"',
             font_size=10, color=DARK_BLUE)

# Turn 3 — 사진 기반 정밀 견적
y += Inches(0.85)
add_text_box(slide, Inches(0.7), y, Inches(4), Inches(0.35),
             "턴 3  사진 기반 정밀 견적 (AI Vision)", font_size=13, color=SAMSUNG_BLUE, bold=True)
y += Inches(0.35)

add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.85), y + Inches(0.08), Inches(5.2), Inches(0.35),
             '👤 [기기 앞면/뒷면 사진 업로드]',
             font_size=11, color=BLACK)

add_rounded_rect(slide, Inches(6.5), y - Inches(0.05), Inches(6.1), Inches(1.85), RGBColor(0xF0, 0xFF, 0xF0), TRADEIN_GREEN)
add_text_box(slide, Inches(6.7), y, Inches(4), Inches(0.3),
             '⚙️ 사진 기반 기기 상태 정밀 분석', font_size=12, color=TRADEIN_GREEN, bold=True)
mcp_t3 = [
    ("입력: 견적ID, Vision 결과 (화면/외관/카메라 상태)", False, MID_GRAY, 10, PP_ALIGN.LEFT),
    ("", False, MID_GRAY, 6, PP_ALIGN.LEFT),
    ("출력:", True, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  • 예상 보상가 범위: 370,000원 ~ 450,000원", False, GREEN, 10, PP_ALIGN.LEFT),
    ("  • 상태 인디케이터: 화면 🟢  바디 🟢  카메라 🟢", False, DARK_GRAY, 10, PP_ALIGN.LEFT),
    ("  ⚠️ 정확한 가격은 전문 업체 수거·분석 후 확정", False, ORANGE, 10, PP_ALIGN.LEFT),
]
add_multi_para(slide, Inches(6.7), y + Inches(0.32), Inches(5.7), Inches(1.4), mcp_t3)

y += Inches(0.6)
add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.55), SKY_BLUE)
add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(5.2), Inches(0.45),
             '🤖 [위젯: 보상가 범위 + 상태 🟢🟢🟢 + disclaimer]\n     "정확한 가격은 전문 업체 수거 및 분석 후 확정됩니다"',
             font_size=10, color=DARK_BLUE)

add_slide_number(slide, 5, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6: Tool List Table
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "MCP 서버 기능(도구) 목록")
add_samsung_footer(slide)

add_text_box(slide, Inches(0.7), Inches(1.25), Inches(11), Inches(0.4),
             "ChatGPT가 대화 중 자동으로 사용하는 9가지 기능", font_size=14, color=MID_BLUE, bold=True)

# Table
from pptx.util import Inches, Pt
table_data = [
    ["#", "기능", "하는 일", "입력 (ChatGPT → 서버)", "출력 (서버 → ChatGPT)", "위젯"],
    ["1", "상담 가이드라인 불러오기", "대화 시작 시 응대 규칙을\nChatGPT에 전달", "(없음)", "응대 톤, 서비스별 대화 흐름,\nVision 사진 가이드", "—"],
    ["2", "Care+ 플랜 조회 및\n가입 가능 여부 확인", "Care+ 보험 플랜 안내,\n구매 시기로 가입 가능 여부 판단", "기기 모델, 구매일", "플랜 비교 (Basic/Premium),\n가입 상태, 수수료, FAQ", "Care+"],
    ["3", "사진 기반 기기 상태 검사\n(Late Enrollment)", "기기 사진을 분석해\nCare+ 늦은 가입 승인/거절", "기기 모델, 구매일,\n사진 분석 결과(화면/외관/카메라)", "가입 가능 여부 (✅/❌),\n항목별 통과 결과, 추천 플랜", "Care+"],
    ["4", "Galaxy Club 플랜별\n가격·혜택 조회", "구독클럽 플랜 3종 소개\n(12/24/36개월)", "보고 싶은 플랜 종류,\nFAQ 포함 여부", "플랜별 월 요금·혜택·잔존가,\n가입 절차, FAQ", "Club"],
    ["5", "구독 vs 일반구매\n비용 비교", "Galaxy Club 구독과\n일반 구매 시 총비용 비교", "비교할 기기 모델,\n플랜 종류", "구독/구매 총비용,\n절약액, 포함 혜택", "비교"],
    ["6", "보상판매 예상 가격 조회", "기기명으로 보상판매\n예상 가격 범위 간편 조회", "기기명, 용량,\n국가, 통신사", "상태별 예상 가격 범위\n(최상~불량), 프로모션", "—"],
    ["7", "보상판매 정식 견적 시작", "상세 정보 입력 후\n정식 보상판매 견적 발급", "기기, 용량, 상태,\n기능·외관 문제, 국가", "견적 번호, 초기 견적가,\n가격 산출 내역", "Trade-in"],
    ["8", "사진 기반 기기 상태\n정밀 분석", "기기 사진을 분석해\n기존 견적 대비 가격 재산정", "견적 번호,\n사진 분석 결과(화면/외관/카메라)", "변경 전/후 가격,\n종합 등급, 다음 절차", "Trade-in"],
    ["9", "보상판매 견적 결과 확인", "이전에 발급받은\n보상판매 견적 결과 조회", "견적 번호", "최종 보상가, 상태,\n다음 절차 안내", "Trade-in"],
]

rows = len(table_data)
cols = len(table_data[0])
tbl_left = Inches(0.5)
tbl_top = Inches(1.75)
tbl_width = Inches(12.3)
tbl_height = Inches(5.0)

table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
table = table_shape.table

# Column widths
col_widths = [Inches(0.35), Inches(2.1), Inches(2.0), Inches(2.5), Inches(2.7), Inches(0.85)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER if c in [0, 5] else PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(9) if r > 0 else Pt(10)
                run.font.name = "맑은 고딕"
                if r == 0:
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                else:
                    run.font.color.rgb = DARK_GRAY
                    if c == 1:
                        run.font.color.rgb = SAMSUNG_BLUE
                        run.font.bold = True
                        run.font.size = Pt(9)

        # Cell fill
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SAMSUNG_BLUE
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ICE_BLUE
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

add_slide_number(slide, 6, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7: Data Inventory — Overview
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "활용 데이터 목록 — 데이터 영역별 개요")
add_samsung_footer(slide)

add_text_box(slide, Inches(0.7), Inches(1.25), Inches(11), Inches(0.4),
             "MCP 서버가 사용하는 4개 데이터 파일과 각각의 역할", font_size=14, color=MID_BLUE, bold=True)

# Data area cards — 2x2 grid
data_cards = [
    ("① 제품 카탈로그", "devices.json", SAMSUNG_BLUE, [
        "43개 Galaxy 기기 정보 (S22~S25, Z4~6, Tab, Watch 등)",
        "기기별 Trade-in 기준가 (용량별 USD 가격)",
        "기기별 정가(MSRP) — Galaxy Club 비용 비교에 사용",
        "스펙 정보 (디스플레이, 프로세서, 카메라, 배터리 등)",
        "출시일, 판매 상태, Trade-in 대상 여부",
    ], [
        "보상판매 예상 가격 조회",
        "보상판매 정식 견적 시작",
        "사진 기반 기기 상태 정밀 분석",
    ]),
    ("② Trade-in 가격 산정 테이블", "devices.json 내 9개 보조 테이블", TRADEIN_GREEN, [
        "상태별 가격 계수 (excellent 100% ~ poor 40%)",
        "6개국 지역별 계수 및 환율 (US/KR/UK/DE/JP/SG)",
        "10개 통신사별 가격 조정 (+$10 ~ -$25)",
        "기능 문제 감가 15종 (배터리, 충전, 카메라 등)",
        "외관 문제 감가 10종 (화면깨짐, 찍힘, 힌지 등)",
        "AI Vision 판정 기준 4부위 (화면/외관/카메라/힌지)",
    ], [
        "보상판매 예상 가격 조회",
        "보상판매 정식 견적 시작",
        "사진 기반 기기 상태 정밀 분석",
    ]),
    ("③ Galaxy Club 구독 플랜", "plans.json", CLUB_PURPLE, [
        "3개 플랜 (12개월/24개월/36개월)",
        "플랜별 기기 가격표 (S26, Z Fold7, Z Flip7 등)",
        "잔존가 보장 금액표 (7건, 모델·용량별)",
        "FAQ 8건, 가입 절차 3단계",
        "구독 라이프사이클 4단계",
    ], [
        "Galaxy Club 플랜별 가격·혜택 조회",
        "구독 vs 일반구매 비용 비교",
    ]),
    ("④ Samsung Care+ 보험", "care_plus.json", CARE_PLUS_BLUE, [
        "2개 플랜 (Basic / Premium)",
        "7개 기기 카테고리별 월 보험료",
        "Late Enrollment 규칙 (60일~1년, Vision 검사)",
        "클레임 프로세스 4단계, 서비스센터 6개국",
        "FAQ 10건, 보장 제외사항 7건",
    ], [
        "Care+ 플랜 조회 및 가입 가능 여부 확인",
        "사진 기반 기기 상태 검사",
    ]),
]

for i, (title, source, accent, items, tools) in enumerate(data_cards):
    col = i % 2
    row = i // 2
    cx = Inches(0.5) + Inches(col * 6.35)
    cy = Inches(1.75) + Inches(row * 2.75)
    card_h = Inches(2.55)

    add_rounded_rect(slide, cx, cy, Inches(6.15), card_h, WHITE, accent)
    # Title bar
    add_rounded_rect(slide, cx, cy, Inches(6.15), Inches(0.5), accent)
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.05), Inches(4), Inches(0.4),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, cx + Inches(4.2), cy + Inches(0.1), Inches(1.8), Inches(0.3),
                 source, font_size=8, color=RGBColor(0xDD, 0xEE, 0xFF), alignment=PP_ALIGN.RIGHT)

    # Data items
    for j, item in enumerate(items):
        add_text_box(slide, cx + Inches(0.15), cy + Inches(0.55 + j * 0.28),
                     Inches(6.0), Inches(0.25),
                     f"• {item}", font_size=9, color=DARK_GRAY)

    # Connected tools label
    tool_y = cy + Inches(0.55 + len(items) * 0.28)
    tools_text = "연결 기능: " + ", ".join(tools)
    add_text_box(slide, cx + Inches(0.15), tool_y, Inches(5.8), Inches(0.25),
                 tools_text, font_size=8, color=accent, bold=True)

# Price formula at bottom
add_rounded_rect(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35), ICE_BLUE, LIGHT_BLUE)
add_text_box(slide, Inches(0.7), Inches(7.08), Inches(11.9), Inches(0.3),
             "가격 산출 공식:  (기준가 x 지역 계수) + 상태 보정 - 기능문제 감가 - 외관문제 감가 + 통신사 조정 + 프로모션 보너스",
             font_size=9, color=MID_BLUE, bold=True)

add_slide_number(slide, 7, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8: Data Inventory — Details
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "활용 데이터 목록 — 상세 현황")
add_samsung_footer(slide)

# Section 1: Device lineup
add_text_box(slide, Inches(0.7), Inches(1.25), Inches(4), Inches(0.35),
             "제품 라인업 (43개 기기)", font_size=14, color=SAMSUNG_BLUE, bold=True)

lineup_data = [
    ["카테고리", "시리즈 / 기기", "수량"],
    ["스마트폰", "Galaxy S25 Ultra ~ S22  (12개)\nGalaxy Z Fold6 ~ Z Flip4  (6개)\nGalaxy A55 / A35 / A25  (3개)", "21"],
    ["태블릿", "Galaxy Tab S10 Ultra ~ Tab S9  (6개)", "6"],
    ["워치", "Galaxy Watch Ultra / Watch7 / Watch6 등", "5"],
    ["이어버즈", "Galaxy Buds3 Pro / Buds3 / Buds2 Pro / Buds FE", "4"],
    ["노트북", "Galaxy Book4 Ultra ~ Book Go  (6개)", "6"],
    ["웨어러블", "Galaxy Ring", "1"],
]

lineup_rows = len(lineup_data)
lineup_cols = 3
lt = slide.shapes.add_table(lineup_rows, lineup_cols, Inches(0.5), Inches(1.65), Inches(6.0), Inches(2.8))
ltbl = lt.table
ltbl.columns[0].width = Inches(1.0)
ltbl.columns[1].width = Inches(4.0)
ltbl.columns[2].width = Inches(0.7)

for r in range(lineup_rows):
    for c in range(lineup_cols):
        cell = ltbl.cell(r, c)
        cell.text = lineup_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if c == 2 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(9) if r > 0 else Pt(10)
                run.font.name = "맑은 고딕"
                run.font.bold = r == 0
                run.font.color.rgb = WHITE if r == 0 else DARK_GRAY
        cell.fill.solid()
        cell.fill.fore_color.rgb = SAMSUNG_BLUE if r == 0 else (ICE_BLUE if r % 2 == 0 else WHITE)

# Section 2: Trade-in sub-tables
add_text_box(slide, Inches(0.7), Inches(4.6), Inches(5), Inches(0.35),
             "Trade-in 가격 산정 보조 테이블 (9종)", font_size=14, color=TRADEIN_GREEN, bold=True)

subtable_data = [
    ["보조 테이블", "내용", "규모"],
    ["상태별 가격 계수", "excellent(100%) / good(85%) / fair(65%) / poor(40%)", "4단계"],
    ["지역별 계수 & 환율", "US / KR / UK / DE / JP / SG — 계수 + 통화 + 환율", "6개국"],
    ["통신사별 조정", "자급제(0) ~ 기타잠금(-$25)", "10종"],
    ["기능 문제 감가", "배터리, 충전, 스피커, 카메라, 지문인식, GPS 등", "15항목"],
    ["외관 문제 감가", "화면깨짐(-$100), 후면깨짐(-$60), 힌지마모(-$40) 등", "10항목"],
    ["AI Vision 판정 기준", "화면/외관/카메라/힌지 각 4단계 → 등급 매핑", "4부위"],
    ["상태 설명 텍스트", "각 상태 등급의 사용자용 설명문", "4건"],
    ["감가상각 일정", "출시 후 연차별 가치 하락률", "5단계"],
    ["프로모션 보너스", "Spring Trade-in Bonus +$50 (스마트폰/태블릿)", "1건"],
]

st_rows = len(subtable_data)
st_cols = 3
st = slide.shapes.add_table(st_rows, st_cols, Inches(0.5), Inches(5.0), Inches(6.0), Inches(2.3))
stbl = st.table
stbl.columns[0].width = Inches(1.5)
stbl.columns[1].width = Inches(3.7)
stbl.columns[2].width = Inches(0.8)

for r in range(st_rows):
    for c in range(st_cols):
        cell = stbl.cell(r, c)
        cell.text = subtable_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if c == 2 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(8) if r > 0 else Pt(9)
                run.font.name = "맑은 고딕"
                run.font.bold = r == 0
                run.font.color.rgb = WHITE if r == 0 else DARK_GRAY
        cell.fill.solid()
        cell.fill.fore_color.rgb = TRADEIN_GREEN if r == 0 else (RGBColor(0xF0, 0xFF, 0xF0) if r % 2 == 0 else WHITE)

# Section 3: Right side — Care+ & Galaxy Club summary
rx = Inches(6.8)

# Galaxy Club data
add_text_box(slide, rx, Inches(1.25), Inches(5.5), Inches(0.35),
             "Galaxy Club 데이터 현황", font_size=13, color=CLUB_PURPLE, bold=True)
gc_items = [
    "플랜 3종: 12개월형(₩6,900/월) / 24개월형(₩6,900/월) / 36개월형(₩8,900/월)",
    "플랜별 대상 기기: S26시리즈, Z Fold7, Z Flip7 (S25는 가입종료)",
    "잔존가 보장: S26 모델·용량별 7건 (예: S26 Ultra 1TB → 12개월 1,272,700원)",
    "FAQ 8건 (반납, 해지, IMEI, 분실보장, 결제 등)",
    "가입 절차 3단계: 구매+신청 → IMEI등록 → 반납+잔존가수령",
    "라이프사이클 4단계: 구매+가입 → IMEI등록 → 기기사용 → 반납",
    "Care+ 플랜별 자동 포함 (12/24mo: 파손, 36mo: 분실/파손)",
    "액세서리 할인 쿠폰 (플랜·기기별 차등)",
]
for j, item in enumerate(gc_items):
    add_text_box(slide, rx + Inches(0.05), Inches(1.6 + j * 0.3), Inches(5.4), Inches(0.28),
                 f"• {item}", font_size=9, color=DARK_GRAY)

# Care+ data
add_text_box(slide, rx, Inches(4.1), Inches(5.5), Inches(0.35),
             "Samsung Care+ 데이터 현황", font_size=13, color=CARE_PLUS_BLUE, bold=True)
cp_items = [
    "플랜 2종: Basic($11.99/월) / Premium($17.99/월) — 플래그십 기준",
    "7개 카테고리별 차등 요금 (스마트폰~노트북, $2.99~$17.99)",
    "자기부담금: Basic 화면$29 vs Premium 화면$0",
    "Late Enrollment: 구매 60일~1년, Vision 검사 통과 시 가입 (수수료 $49)",
    "Vision 통과 기준: 화면·외관·카메라 각각 \"양호\" 이상",
    "클레임 프로세스 4단계 (접수→확인→수리/교체→완료)",
    "서비스센터: 전 세계 3,200개 (US 420, KR 180 등 6개국)",
    "FAQ 10건, 보장 제외사항 7건",
]
for j, item in enumerate(cp_items):
    add_text_box(slide, rx + Inches(0.05), Inches(4.45 + j * 0.3), Inches(5.4), Inches(0.28),
                 f"• {item}", font_size=9, color=DARK_GRAY)

# Note at bottom right
add_rounded_rect(slide, rx, Inches(6.95), Inches(5.8), Inches(0.45), RGBColor(0xFF, 0xF5, 0xEE), ORANGE)
add_text_box(slide, rx + Inches(0.15), Inches(7.0), Inches(5.5), Inches(0.35),
             "참고: promotions.json (22개 캠페인)은 현재 서버에서 미사용 — 향후 프로모션 기능 확장 시 활용 예정",
             font_size=9, color=ORANGE, bold=True)

add_slide_number(slide, 8, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9: Widget Overview
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "위젯 (사용자에게 보이는 화면) & AI Vision 활용")
add_samsung_footer(slide)

# Widget cards
widget_cards = [
    ("Care+ 위젯", CARE_PLUS_BLUE, [
        "Basic / Premium 플랜 비교 표",
        "가입 상태 배너 (즉시/Late/만료)",
        "Vision 검사 결과 인디케이터 🟢🟡🔴",
        "CTA: \"Care+ 가입하기\" 버튼",
    ]),
    ("Galaxy Club 위젯", CLUB_PURPLE, [
        "12/24/36개월 플랜 탭",
        "월 요금, 잔존가, 혜택 상세",
        "구독 라이프사이클 타임라인",
        "가입 절차 단계 표시",
    ]),
    ("비교 위젯", MID_BLUE, [
        "구독 vs 일반구매 비용 비교 차트",
        "절약액 강조 배너",
        "포함 혜택 목록",
        "Trade-in 정보 (해당 시)",
    ]),
    ("Trade-in 위젯", TRADEIN_GREEN, [
        "기기 검색 입력창",
        "상태별 가격 범위 표시",
        "견적 결과 (Before / After)",
        "상태 인디케이터 + CTA 버튼",
    ]),
]

for i, (name, accent, items) in enumerate(widget_cards):
    cx = Inches(0.5) + Inches(i * 3.15)
    cy = Inches(1.3)
    ch = Inches(2.7)
    add_rounded_rect(slide, cx, cy, Inches(2.95), ch, WHITE, accent)
    # Title bar
    add_rounded_rect(slide, cx, cy, Inches(2.95), Inches(0.55), accent)
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.1), Inches(2.65), Inches(0.4),
                 name, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text_box(slide, cx + Inches(0.15), cy + Inches(0.7 + j * 0.42), Inches(2.65), Inches(0.4),
                     f"• {item}", font_size=10, color=DARK_GRAY)

# AI Vision section
vy = Inches(4.3)
add_rounded_rect(slide, Inches(0.5), vy, Inches(12.3), Inches(2.5), RGBColor(0xFD, 0xF5, 0xE6), ORANGE)
add_text_box(slide, Inches(0.8), vy + Inches(0.12), Inches(11.5), Inches(0.4),
             "🔍 핵심 기술: AI Vision 활용", font_size=16, color=ORANGE, bold=True)

add_text_box(slide, Inches(0.8), vy + Inches(0.55), Inches(11.5), Inches(0.4),
             "3개 시나리오 중 2개에서 ChatGPT의 Vision(이미지 분석) 기능을 활용합니다.",
             font_size=12, color=DARK_GRAY)

# Two Vision use-cases side by side
vision_cases = [
    ("Care+ Late Enrollment", CARE_PLUS_BLUE, "사진 업로드 → ChatGPT Vision 분석\n→ 화면/외관/카메라 상태 판정\n→ MCP 서버에서 가입 승인/거절"),
    ("Trade-in 정밀 견적", TRADEIN_GREEN, "사진 업로드 → ChatGPT Vision 분석\n→ 화면/외관/카메라 상태 판정\n→ MCP 서버에서 가격 재산정"),
]
for i, (title, clr, desc) in enumerate(vision_cases):
    vx = Inches(0.8) + Inches(i * 6.2)
    add_rounded_rect(slide, vx, vy + Inches(1.0), Inches(5.7), Inches(1.3), WHITE, clr)
    add_text_box(slide, vx + Inches(0.15), vy + Inches(1.08), Inches(5.4), Inches(0.35),
                 title, font_size=12, color=clr, bold=True)
    add_text_box(slide, vx + Inches(0.15), vy + Inches(1.4), Inches(5.4), Inches(0.8),
                 desc, font_size=10, color=DARK_GRAY)

add_slide_number(slide, 9, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10: Tech Requirements — PoC / MVP
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "기술 요건 — PoC / MVP 수준")
add_samsung_footer(slide)

add_rounded_rect(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45), ICE_BLUE, LIGHT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.28), Inches(11.5), Inches(0.35),
             "전제: 내부 시연 또는 제한된 사용자 대상 검증 목적",
             font_size=12, color=MID_BLUE, bold=True)

# 3 category cards
poc_cards = [
    ("서버 & 인프라", SAMSUNG_BLUE, [
        "Node.js 서버 1대 (현재 PoC 구조)",
        "클라우드 VM 또는 컨테이너 1개",
        "HTTPS 인증서 (Let's Encrypt)",
        "고정 도메인 1개 (ngrok 대체)",
    ]),
    ("데이터", MID_BLUE, [
        "JSON 파일 또는 간단한 DB",
        "(현재 mock 데이터를 실 데이터로 교체)",
        "가격 정보 수동 업데이트",
        "사용자 데이터 저장 불필요 (세션 단위)",
    ]),
    ("ChatGPT 연동", LIGHT_BLUE, [
        "OpenAI Apps 파트너 등록 신청",
        "Streamable HTTP Transport (현재 구현)",
        "Widget: HTML/CSS/JS (정적 iframe)",
        "MCP 프로토콜 표준 준수",
    ]),
]

for i, (title, clr, items) in enumerate(poc_cards):
    cx = Inches(0.5) + Inches(i * 4.15)
    cy = Inches(1.9)
    add_rounded_rect(slide, cx, cy, Inches(3.95), Inches(2.8), WHITE, clr)
    add_rounded_rect(slide, cx, cy, Inches(3.95), Inches(0.5), clr)
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.08), Inches(3.65), Inches(0.35),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text_box(slide, cx + Inches(0.15), cy + Inches(0.6 + j * 0.45), Inches(3.65), Inches(0.4),
                     f"• {item}", font_size=11, color=DARK_GRAY)

# Team & timeline
ty = Inches(5.0)
add_rounded_rect(slide, Inches(0.5), ty, Inches(12.3), Inches(1.8), ICE_BLUE, SAMSUNG_BLUE)
add_text_box(slide, Inches(0.8), ty + Inches(0.1), Inches(3), Inches(0.4),
             "예상 인력 & 기간", font_size=14, color=SAMSUNG_BLUE, bold=True)

team_data = [
    ["백엔드 개발 1명", "MCP 서버 개발, 데이터 연동"],
    ["프론트엔드 개발 1명", "위젯 UI 개발"],
    ["기획/데이터 1명", "서비스 시나리오, 데이터 준비"],
]
for i, (role, work) in enumerate(team_data):
    tx = Inches(0.8) + Inches(i * 3.5)
    add_text_box(slide, tx, ty + Inches(0.5), Inches(3.3), Inches(0.35),
                 role, font_size=12, color=SAMSUNG_BLUE, bold=True)
    add_text_box(slide, tx, ty + Inches(0.85), Inches(3.3), Inches(0.35),
                 work, font_size=10, color=DARK_GRAY)

# Timeline highlight
add_rounded_rect(slide, Inches(0.8), ty + Inches(1.25), Inches(2.5), Inches(0.45), SAMSUNG_BLUE)
add_text_box(slide, Inches(0.95), ty + Inches(1.3), Inches(2.3), Inches(0.35),
             "총 예상 기간: 4~6주", font_size=13, color=WHITE, bold=True)

add_slide_number(slide, 10, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11: Production — Samsung Account & Pre-filled Redirect
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "Production — 삼성 계정 연동 & 서비스 페이지 연결")
add_samsung_footer(slide)

add_rounded_rect(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45), ICE_BLUE, LIGHT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.28), Inches(11.5), Inches(0.35),
             "대화에서 수집한 정보로 삼성닷컴 신청 페이지에 바로 연결하려면?",
             font_size=12, color=MID_BLUE, bold=True)

# Left: Samsung Account Integration
add_rounded_rect(slide, Inches(0.5), Inches(1.85), Inches(6.15), Inches(2.65), WHITE, SAMSUNG_BLUE)
add_rounded_rect(slide, Inches(0.5), Inches(1.85), Inches(6.15), Inches(0.45), SAMSUNG_BLUE)
add_text_box(slide, Inches(0.65), Inches(1.9), Inches(5.85), Inches(0.35),
             "Samsung Account 연동", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
sa_items = [
    "Samsung Account는 OAuth 2.0 + OpenID Connect 표준 사용",
    "client_id 발급에 삼성 파트너 계약 필요 (자체 등록 불가)",
    "Samsung Account 인프라: AWS EKS, 70+ 마이크로서비스, 4개 리전",
    "연동 흐름: ChatGPT 위젯 CTA 클릭 → Samsung Account 로그인",
    "   → OAuth state에 대화 컨텍스트(견적ID 등) 포함 → 콜백 후 신청 페이지 연결",
    "OpenAI Apps SDK가 OAuth 2.1 흐름 지원 → state 파라미터로 컨텍스트 전달",
]
for j, item in enumerate(sa_items):
    add_text_box(slide, Inches(0.65), Inches(2.4 + j * 0.3), Inches(5.85), Inches(0.28),
                 f"• {item}", font_size=9, color=DARK_GRAY)

# Right: Pre-filled Redirect Flow
add_rounded_rect(slide, Inches(6.85), Inches(1.85), Inches(6.15), Inches(2.65), WHITE, MID_BLUE)
add_rounded_rect(slide, Inches(6.85), Inches(1.85), Inches(6.15), Inches(0.45), MID_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.85), Inches(0.35),
             "대화 → 신청 페이지 사전입력 연결", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
rd_items = [
    "패턴: 서버에서 단기 토큰(5~15분) 발급 → URL에 토큰만 포함",
    "   → 랜딩페이지가 토큰으로 서버 콜백 → 저장된 대화 정보 수신 → 폼 사전입력",
    "전달 가능 정보: 기기 모델, 용량, 상태, 견적ID, 보상가, 플랜 선택 등",
    "보안: 토큰 1회 사용 후 무효화, HMAC 서명으로 위변조 방지",
    "PII(IMEI, 이름, 연락처)는 URL 파라미터에 절대 포함 금지",
    "랜딩페이지: 삼성닷컴 SAP Commerce 기반 → 토큰 수신 API 개발 필요",
]
for j, item in enumerate(rd_items):
    add_text_box(slide, Inches(7.0), Inches(2.4 + j * 0.3), Inches(5.85), Inches(0.28),
                 f"• {item}", font_size=9, color=DARK_GRAY)

# Flow diagram
add_text_box(slide, Inches(0.7), Inches(4.7), Inches(4), Inches(0.35),
             "예시 흐름: Trade-in 보상판매 신청", font_size=12, color=SAMSUNG_BLUE, bold=True)

flow_steps = [
    ("① 대화에서\n견적 완료", RGBColor(0x4A, 0x4A, 0x4A)),
    ("② 위젯 CTA\n\"신청하기\" 클릭", SAMSUNG_BLUE),
    ("③ Samsung Account\nOAuth 로그인", MID_BLUE),
    ("④ MCP 서버에서\n토큰 발급", CARE_PLUS_BLUE),
    ("⑤ 삼성닷컴 신청 페이지\n(폼 사전입력 완료)", TRADEIN_GREEN),
]
for i, (text, clr) in enumerate(flow_steps):
    fx = Inches(0.5) + Inches(i * 2.55)
    add_rounded_rect(slide, fx, Inches(5.1), Inches(2.2), Inches(0.9), clr)
    add_text_box(slide, fx + Inches(0.1), Inches(5.15), Inches(2.0), Inches(0.8),
                 text, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, fx + Inches(2.1), Inches(5.3), Inches(0.5), Inches(0.5),
                     "→", font_size=18, color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom notes
add_rounded_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), RGBColor(0xFF, 0xF5, 0xEE), ORANGE)
add_text_box(slide, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.6),
             "삼성전자 준비 사항:  ① Samsung Account OAuth client_id 발급 (파트너 계약)  "
             "② 삼성닷컴 랜딩페이지에 토큰 수신 API 개발 (SAP Commerce 확장)  "
             "③ 위젯 CTA → 삼성닷컴 딥링크 URL 스키마 정의",
             font_size=10, color=ORANGE, bold=True)

add_slide_number(slide, 11, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 12: Production — Data Sourcing from Samsung Systems
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "Production — 실 데이터 소싱 & 삼성 내부 시스템 연동")
add_samsung_footer(slide)

add_rounded_rect(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45), ICE_BLUE, LIGHT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.28), Inches(11.5), Inches(0.35),
             "현재 Mock 데이터 → 삼성 내부 시스템의 실제 데이터로 전환",
             font_size=12, color=MID_BLUE, bold=True)

# Data sourcing table
ds_data = [
    ["현재 Mock 데이터", "실 데이터 소스 (삼성 시스템)", "연동 방식", "갱신 주기"],
    ["제품 카탈로그\n(devices.json, 43개 기기)", "SAP Commerce Cloud (Hybris)\n삼성닷컴 PIM (제품정보관리)", "SAP Commerce OCC API\n또는 Kafka 스트림 구독", "기기 출시 시\n(분기 1~2회)"],
    ["기기 정가 (MSRP)\n용량별 가격", "SAP Commerce 가격 엔진\nKafka 실시간 가격 스트림", "REST API 조회\n또는 이벤트 수신", "가격 변동 시\n(수시)"],
    ["Trade-in 기준가\n상태별 계수", "삼성닷컴 Trade-in 룩업 테이블\n(SAP Commerce 프로모션 룰)", "내부 API 조회", "모델 출시/분기별\n가격 재산정 시"],
    ["Trade-in 실물 감정\n(한국)", "민팃 (SK Networks)\nATM 6,600대, AI 31포인트 검사", "민팃 파트너 API\n(현재 존재 여부 확인 필요)", "실시간\n(ATM 감정 결과)"],
    ["Care+ 플랜/가입/클레임", "Servify 플랫폼\n(integrations.servify.com)", "Servify REST API\n(IMEI 조회, 플랜, 가입 적격)", "실시간"],
    ["Galaxy Club 구독", "Samsung Account + SAP Commerce\n(IMEI 기반 구독 관리)", "Samsung Account API\nSAP Commerce 구독 모듈", "실시간"],
    ["고객 인증", "Samsung Account\n(AWS EKS, 4개 리전, 2.7M RPS)", "OAuth 2.0 + OIDC\n(client_id 파트너 발급)", "실시간"],
]

ds_rows = len(ds_data)
ds_cols = 4
dt = slide.shapes.add_table(ds_rows, ds_cols, Inches(0.3), Inches(1.85), Inches(12.7), Inches(4.6))
dtbl = dt.table
dtbl.columns[0].width = Inches(2.4)
dtbl.columns[1].width = Inches(3.8)
dtbl.columns[2].width = Inches(3.2)
dtbl.columns[3].width = Inches(1.5)

for r in range(ds_rows):
    for c in range(ds_cols):
        cell = dtbl.cell(r, c)
        cell.text = ds_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(8) if r > 0 else Pt(9)
                run.font.name = "맑은 고딕"
                run.font.bold = r == 0
                run.font.color.rgb = WHITE if r == 0 else DARK_GRAY
        cell.fill.solid()
        cell.fill.fore_color.rgb = SAMSUNG_BLUE if r == 0 else (ICE_BLUE if r % 2 == 0 else WHITE)

# Bottom: Samsung action items
add_rounded_rect(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.8), RGBColor(0xFF, 0xF5, 0xEE), ORANGE)
add_text_box(slide, Inches(0.5), Inches(6.62), Inches(12.3), Inches(0.75),
             "삼성전자 준비 사항:\n"
             "① SAP Commerce OCC API 접근 권한 부여 (제품/가격 데이터)  "
             "② Servify 파트너 API 연동 계약 (Care+ 데이터)  "
             "③ 민팃 API 연동 가능 여부 확인 (Trade-in 실물 감정)  "
             "④ MCP 서버 전용 서비스 계정 발급 (API Gateway 경유, rate limit 설정)",
             font_size=9, color=ORANGE, bold=True)

add_slide_number(slide, 12, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 13: Production — Security & Compliance
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WHITE)
add_header_bar(slide, "Production — 보안, 개인정보, 인프라 요건")
add_samsung_footer(slide)

# 3 columns
sec_cards = [
    ("개인정보보호법 (PIPA) 준수", DARK_BLUE, [
        "2023년 개정 PIPA 적용 (시행 2023.9.15)",
        "ChatGPT(OpenAI)로의 국외 이전:",
        "  → 이용자 동의 또는 계약 위탁 고지 필요",
        "72시간 이내 침해 통지 의무",
        "자동화된 결정(AI 견적)에 대한 거부권 보장",
        "IMEI, 구매이력 등 수집 목적 명시 필수",
        "대화 내 PII 최소 보관, 목적 달성 후 파기",
        "개인정보 처리방침 + DPA(OpenAI) 체결",
    ]),
    ("보안 아키텍처", SAMSUNG_BLUE, [
        "MCP 서버: Samsung Cloud 또는 AWS EKS 배포",
        "삼성 내부 네트워크 ↔ MCP: API Gateway 경유",
        "  → mTLS + API Key + IP 화이트리스트",
        "ChatGPT ↔ MCP: HTTPS + OAuth 2.0 토큰",
        "위젯 CSP: 삼성 도메인만 허용 (allowlist)",
        "Samsung Account OAuth: PKCE 필수",
        "모든 도구 호출 감사 로그 (who/when/what)",
        "Knox 보안 인증 기준 적용 (CC EAL4+, FIPS)",
    ]),
    ("인프라 & 운영", MID_BLUE, [
        "컨테이너 오케스트레이션 (K8s), Auto-scaling",
        "Multi-AZ 배포, SLA 99.9%, 응답 <2초",
        "APM 모니터링 + 알림 (Datadog/CloudWatch)",
        "CI/CD: 자동 빌드·배포·롤백",
        "E2E 테스트 자동화 (시나리오별)",
        "Vision 정확도 벤치마크 및 정기 검증",
        "다국어 위젯 (한/영/일/독)",
        "인간 상담원 에스컬레이션 프로세스",
    ]),
]

for i, (title, clr, items) in enumerate(sec_cards):
    cx = Inches(0.3) + Inches(i * 4.3)
    cy = Inches(1.3)
    add_rounded_rect(slide, cx, cy, Inches(4.1), Inches(4.5), WHITE, clr)
    add_rounded_rect(slide, cx, cy, Inches(4.1), Inches(0.45), clr)
    add_text_box(slide, cx + Inches(0.1), cy + Inches(0.05), Inches(3.9), Inches(0.35),
                 title, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text_box(slide, cx + Inches(0.1), cy + Inches(0.55 + j * 0.45),
                     Inches(3.9), Inches(0.42),
                     f"{'  ' if item.startswith('  ') else '• '}{item.strip()}", font_size=8, color=DARK_GRAY)

# Bottom: team & timeline
add_rounded_rect(slide, Inches(0.3), Inches(5.95), Inches(12.7), Inches(0.95), ICE_BLUE, SAMSUNG_BLUE)
add_text_box(slide, Inches(0.5), Inches(6.0), Inches(6), Inches(0.35),
             "예상 인력 & 기간 (Production)", font_size=12, color=SAMSUNG_BLUE, bold=True)

prod_team = [
    ("백엔드 2~3명", "MCP 서버, SAP/Servify/민팃 API 연동, 인증"),
    ("프론트 1~2명", "위젯 UI/UX, 다국어, 접근성"),
    ("AI/ML 1명", "Vision 정확도 검증, 벤치마크"),
    ("DevOps 1명", "K8s, CI/CD, 모니터링, 보안"),
    ("기획+법무 1~2명", "시나리오, 파트너 계약, PIPA 검토"),
]
for i, (role, desc) in enumerate(prod_team):
    tx = Inches(0.5) + Inches(i * 2.5)
    add_text_box(slide, tx, Inches(6.35), Inches(2.4), Inches(0.25),
                 role, font_size=9, color=SAMSUNG_BLUE, bold=True)
    add_text_box(slide, tx, Inches(6.58), Inches(2.4), Inches(0.25),
                 desc, font_size=7, color=DARK_GRAY)

add_rounded_rect(slide, Inches(0.5), Inches(6.9), Inches(2.8), Inches(0.4), SAMSUNG_BLUE)
add_text_box(slide, Inches(0.65), Inches(6.93), Inches(2.5), Inches(0.35),
             "총 예상 기간: 3~5개월", font_size=12, color=WHITE, bold=True)

add_slide_number(slide, 13, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 14: Summary
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, DARK_BLUE)
add_bg_rect(slide, MID_BLUE, Inches(0), Inches(2.8), W, Inches(0.04))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.5),
             "SAMSUNG", font_size=18, color=RGBColor(0x88, 0xAA, 0xDD), bold=True, font_name="Arial")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "요약 — 이 프로젝트가 보여주는 것", font_size=32, color=WHITE, bold=True)

summary_items = [
    ("ChatGPT 안에서 삼성 서비스를 자연어로 이용",
     "별도 앱 설치 없이, 대화만으로 Care+ 가입, 구독 비교, 보상판매 견적까지"),
    ("AI Vision을 활용한 기기 상태 자동 판정",
     "사진만 올리면 화면/외관/카메라 상태를 분석하여 가격 산정 또는 가입 승인"),
    ("인터랙티브 위젯",
     "비교표, 차트, 버튼이 있는 화면을 ChatGPT 대화 안에서 직접 표시"),
    ("MCP 표준 프로토콜 기반",
     "ChatGPT뿐 아니라 MCP 지원 AI 어시스턴트에도 동일 서버 연결 가능"),
]

nums = ["01", "02", "03", "04"]
colors_accent = [CARE_PLUS_BLUE, LIGHT_BLUE, CLUB_PURPLE, TRADEIN_GREEN]

for i, ((title, desc), num, accent) in enumerate(zip(summary_items, nums, colors_accent)):
    sy = Inches(3.2) + Inches(i * 1.0)
    # Number circle
    add_rounded_rect(slide, Inches(0.8), sy, Inches(0.7), Inches(0.7), accent)
    add_text_box(slide, Inches(0.8), sy + Inches(0.1), Inches(0.7), Inches(0.5),
                 num, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Arial")
    # Title and desc
    add_text_box(slide, Inches(1.8), sy + Inches(0.05), Inches(10), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.8), sy + Inches(0.45), Inches(10), Inches(0.4),
                 desc, font_size=12, color=SKY_BLUE)

add_slide_number(slide, 14, TOTAL_SLIDES)


# ── Save ──
output_path = os.path.join(r"C:\Users\user\samsung-galaxy-mcp-share\samsung-galaxy-mcp",
                           "Samsung_Galaxy_MCP_소개.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
